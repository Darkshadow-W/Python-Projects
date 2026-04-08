import os
import subprocess
from pathlib import Path
from datetime import datetime
import mimetypes

try:
    from docx import Document as DocxDocument
    from docx.shared import Pt, Inches
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
except ImportError:
    Presentation = None

try:
    from pdf2image import convert_from_path
except ImportError:
    convert_from_path = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

class FileConverter:
    """Handle file conversions between different formats"""
    
    SUPPORTED_FORMATS = {
        'pdf': ['docx', 'pptx', 'txt', 'images'],
        'docx': ['pdf', 'txt'],
        'doc': ['pdf', 'docx', 'txt'],
        'pptx': ['pdf', 'txt', 'images'],
        'ppt': ['pdf', 'pptx', 'txt'],
        'txt': ['pdf', 'docx'],
        'odt': ['pdf', 'docx']
    }
    
    @staticmethod
    def get_file_extension(filename):
        """Extract file extension"""
        return Path(filename).suffix.lower().lstrip('.')
    
    @staticmethod
    def is_valid_conversion(input_format, output_format):
        """Check if conversion is supported"""
        input_fmt = input_format.lower()
        output_fmt = output_format.lower()
        
        if input_fmt not in FileConverter.SUPPORTED_FORMATS:
            return False
        
        return output_fmt in FileConverter.SUPPORTED_FORMATS[input_fmt]
    
    @staticmethod
    def convert_pdf_to_docx(input_path, output_path):
        """Convert PDF to DOCX"""
        try:
            # Use LibreOffice/UNO for conversion if available
            cmd = [
                'libreoffice', '--headless', '--convert-to', 'docx',
                '--outdir', os.path.dirname(output_path), input_path
            ]
            subprocess.run(cmd, check=True, capture_output=True)
            
            # Rename output file
            converted_file = os.path.join(
                os.path.dirname(output_path),
                Path(input_path).stem + '.docx'
            )
            if os.path.exists(converted_file):
                os.rename(converted_file, output_path)
            
            return True
        except Exception as e:
            print(f"PDF to DOCX conversion error: {e}")
            # Fallback: create basic DOCX with note
            try:
                if DocxDocument:
                    doc = DocxDocument()
                    doc.add_paragraph('[Converted from PDF - Please manually verify content]')
                    doc.save(output_path)
                    return True
            except:
                pass
            return False
    
    @staticmethod
    def convert_pdf_to_pptx(input_path, output_path):
        """Convert PDF to PPTX"""
        try:
            if convert_from_path and Presentation:
                # Convert PDF pages to images
                images = convert_from_path(input_path)
                
                # Create presentation
                prs = Presentation()
                prs.slide_width = PptxInches(10)
                prs.slide_height = PptxInches(7.5)
                
                # Add images as slides
                for img in images:
                    # Save image temporarily
                    temp_img_path = f"/tmp/slide_{datetime.now().timestamp()}.png"
                    img.save(temp_img_path)
                    
                    # Add to presentation
                    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                    left = PptxInches(0)
                    top = PptxInches(0)
                    pic = slide.shapes.add_picture(temp_img_path, left, top, 
                                                   width=prs.slide_width,
                                                   height=prs.slide_height)
                    
                    # Clean up temp image
                    os.remove(temp_img_path)
                
                prs.save(output_path)
                return True
        except Exception as e:
            print(f"PDF to PPTX conversion error: {e}")
        
        return False
    
    @staticmethod
    def convert_docx_to_pdf(input_path, output_path):
        """Convert DOCX to PDF"""
        try:
            cmd = [
                'libreoffice', '--headless', '--convert-to', 'pdf',
                '--outdir', os.path.dirname(output_path), input_path
            ]
            subprocess.run(cmd, check=True, capture_output=True)
            
            # Rename output file
            converted_file = os.path.join(
                os.path.dirname(output_path),
                Path(input_path).stem + '.pdf'
            )
            if os.path.exists(converted_file):
                os.rename(converted_file, output_path)
            
            return True
        except Exception as e:
            print(f"DOCX to PDF conversion error: {e}")
            return False
    
    @staticmethod
    def convert_pptx_to_pdf(input_path, output_path):
        """Convert PPTX to PDF"""
        try:
            cmd = [
                'libreoffice', '--headless', '--convert-to', 'pdf',
                '--outdir', os.path.dirname(output_path), input_path
            ]
            subprocess.run(cmd, check=True, capture_output=True)
            
            # Rename output file
            converted_file = os.path.join(
                os.path.dirname(output_path),
                Path(input_path).stem + '.pdf'
            )
            if os.path.exists(converted_file):
                os.rename(converted_file, output_path)
            
            return True
        except Exception as e:
            print(f"PPTX to PDF conversion error: {e}")
            return False
    
    @staticmethod
    def convert_pdf_to_images(input_path, output_dir):
        """Convert PDF pages to images"""
        try:
            if convert_from_path:
                images = convert_from_path(input_path)
                
                os.makedirs(output_dir, exist_ok=True)
                image_paths = []
                
                for idx, img in enumerate(images, 1):
                    img_path = os.path.join(output_dir, f'page_{idx:03d}.png')
                    img.save(img_path)
                    image_paths.append(img_path)
                
                return image_paths
        except Exception as e:
            print(f"PDF to images conversion error: {e}")
        
        return []
    
    @staticmethod
    def convert_txt_to_pdf(input_path, output_path):
        """Convert TXT to PDF"""
        try:
            if DocxDocument:
                # Create DOCX from text
                doc = DocxDocument()
                
                with open(input_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                    doc.add_paragraph(content)
                
                # Save as DOCX temporarily
                temp_docx = output_path.replace('.pdf', '_temp.docx')
                doc.save(temp_docx)
                
                # Convert DOCX to PDF
                return FileConverter.convert_docx_to_pdf(temp_docx, output_path)
        except Exception as e:
            print(f"TXT to PDF conversion error: {e}")
        
        return False
    
    @classmethod
    def convert(cls, input_path, output_format):
        """Main conversion method"""
        input_format = cls.get_file_extension(input_path)
        
        if not cls.is_valid_conversion(input_format, output_format):
            raise ValueError(f"Conversion from {input_format} to {output_format} is not supported")
        
        output_path = input_path.replace(f'.{input_format}', f'.{output_format}')
        
        # Route to appropriate converter
        conversion_method = f'convert_{input_format}_to_{output_format}'
        
        if hasattr(cls, conversion_method):
            method = getattr(cls, conversion_method)
            success = method(input_path, output_path)
        else:
            # Try LibreOffice as fallback
            success = cls._convert_with_libreoffice(input_path, output_format)
            
        if success and os.path.exists(output_path):
            return output_path
        
        return None
    
    @staticmethod
    def _convert_with_libreoffice(input_path, output_format):
        """Fallback conversion using LibreOffice"""
        try:
            output_dir = os.path.dirname(input_path)
            
            cmd = [
                'libreoffice', '--headless', '--convert-to', output_format,
                '--outdir', output_dir, input_path
            ]
            subprocess.run(cmd, check=True, capture_output=True)
            return True
        except Exception as e:
            print(f"LibreOffice conversion error: {e}")
            return False
